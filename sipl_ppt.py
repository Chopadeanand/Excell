"""
sipl_ppt.py  —  HiRATE Report → PowerPoint generator (python-pptx)
===================================================================
Reads one or more *_REPORT.xlsx files (output of sipl_report.py),
extracts chart + summary data, and builds 2 slides per project into
a single cumulative HiRATE_Report.pptx:

  Slide 1 (per project):  Division-wise clustered bar chart (full width)
  Slide 2 (per project):  Category-wise clustered bar chart + Observations pie

Dependencies:  pip install python-pptx openpyxl
"""

import os, sys, glob, tempfile, argparse, re
import openpyxl

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.chart import XL_CHART_TYPE, XL_LEGEND_POSITION, XL_LABEL_POSITION
from pptx.chart.data import ChartData

DEFAULT_OUTPUT = "HiRATE_Report.pptx"
SCRIPT_DIR     = os.path.dirname(os.path.abspath(__file__))

# ── Colours ───────────────────────────────────────────────────────────────────
NAVY        = RGBColor(0x1F, 0x4E, 0x79)
ORANGE      = RGBColor(0xED, 0x7D, 0x31)
GREEN       = RGBColor(0x92, 0xD0, 0x50)
WHITE       = RGBColor(0xFF, 0xFF, 0xFF)
GRAY        = RGBColor(0x55, 0x55, 0x55)
BG_HDR      = RGBColor(0x1A, 0x2E, 0x45)
BG_SLD      = RGBColor(0xF7, 0xF9, 0xFC)
BLUE        = RGBColor(0x00, 0x70, 0xC0)
DARK        = RGBColor(0x0D, 0x1B, 0x2A)
CYAN        = RGBColor(0x00, 0xB0, 0xF0)
GREEN_DARK  = RGBColor(0x2E, 0x7D, 0x32)
ORANGE_DARK = RGBColor(0xE6, 0x51, 0x00)
GREEN_BG    = RGBColor(0xE8, 0xF5, 0xE1)
ORANGE_BG   = RGBColor(0xFF, 0xF3, 0xE0)

# ── Slide size: 16:9 ─────────────────────────────────────────────────────────
SW = Inches(10)
SH = Inches(5.625)


# ══════════════════════════════════════════════════════════════════════════════
# LOW-LEVEL HELPERS
# ══════════════════════════════════════════════════════════════════════════════

def add_rect(slide, x, y, w, h, fill_rgb, line_rgb=None, line_pt=0):
    from pptx.enum.shapes import MSO_SHAPE_TYPE
    shape = slide.shapes.add_shape(
        1,  # MSO_CONNECTOR_TYPE.STRAIGHT / rectangle = 1
        Inches(x), Inches(y), Inches(w), Inches(h)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_rgb
    if line_rgb:
        shape.line.color.rgb = line_rgb
        shape.line.width = Pt(line_pt)
    else:
        shape.line.fill.background()
    return shape


def add_textbox(slide, x, y, w, h, text, size, bold=False, italic=False,
                colour=WHITE, align=PP_ALIGN.LEFT, face="Calibri"):
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text            = text
    run.font.size       = Pt(size)
    run.font.bold       = bold
    run.font.italic     = italic
    run.font.color.rgb  = colour
    run.font.name       = face
    return txb


def add_stat_box(slide, x, y, w, h, number, subtitle, bg, border, num_col):
    add_rect(slide, x, y, w, h, bg, border, 1.5)
    txb = slide.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf  = txb.text_frame
    tf.word_wrap = True
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text           = str(number)
    r1.font.size      = Pt(18)
    r1.font.bold      = True
    r1.font.color.rgb = num_col
    r1.font.name      = "Calibri"
    p2 = tf.add_paragraph()
    p2.alignment = PP_ALIGN.CENTER
    r2 = p2.add_run()
    r2.text           = subtitle
    r2.font.size      = Pt(7)
    r2.font.color.rgb = GRAY
    r2.font.name      = "Calibri"


# ══════════════════════════════════════════════════════════════════════════════
# HEADER / FOOTER
# ══════════════════════════════════════════════════════════════════════════════

def add_header(slide, title, proj):
    add_rect(slide, 0, 0, 10, 0.72, BG_HDR)
    add_textbox(slide, 0.18, 0.08, 1.4, 0.56, "HiRATE",
                20, bold=True, colour=CYAN)
    add_textbox(slide, 1.7, 0.08, 6.8, 0.56, title,
                13, bold=True, colour=WHITE)
    add_rect(slide, 8.55, 0.13, 1.3, 0.46, BLUE)
    add_textbox(slide, 8.55, 0.13, 1.3, 0.46, proj,
                10, bold=True, colour=WHITE, align=PP_ALIGN.CENTER)


def add_footer(slide, label):
    add_rect(slide, 0, 5.42, 10, 0.2, DARK)
    add_textbox(slide, 0.2, 5.43, 7, 0.18,
                "HiRATE Audit Report  |  Confidential",
                7, colour=RGBColor(0x88, 0x88, 0x88))
    add_textbox(slide, 8.5, 5.43, 1.4, 0.18, label,
                7, colour=RGBColor(0x88, 0x88, 0x88), align=PP_ALIGN.RIGHT)


# ══════════════════════════════════════════════════════════════════════════════
# CHART HELPERS
# ══════════════════════════════════════════════════════════════════════════════

def _hex(rgb): return '%02X%02X%02X' % (rgb[0], rgb[1], rgb[2])


def _set_series_colour(series, rgb):
    from pptx.oxml.ns import qn
    from lxml import etree
    sp = series._element.get_or_add_spPr()
    for old in sp.findall(qn('a:solidFill')):
        sp.remove(old)
    sf   = etree.SubElement(sp, qn('a:solidFill'))
    srgb = etree.SubElement(sf, qn('a:srgbClr'))
    srgb.set('val', _hex(rgb))


def add_bar_chart(slide, x, y, w, h, categories, series_list, y_max=None, lbl_size=8):
    """series_list = [("Name", [values], RGBColor), ...]"""
    cd = ChartData()
    cd.categories = categories
    for name, vals, _ in series_list:
        cd.add_series(name, vals)

    gf = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(x), Inches(y), Inches(w), Inches(h), cd
    )
    chart = gf.chart

    for i, (_, _, col) in enumerate(series_list):
        _set_series_colour(chart.series[i], col)
        dl = chart.series[i].data_labels
        dl.showValue = True
        dl.position = XL_LABEL_POSITION.OUTSIDE_END
        dl.font.size = Pt(lbl_size)
        dl.font.bold = True
        dl.font.color.rgb = col
        # python-pptx bug: showValue at series level stays 0 — fix via lxml
        from lxml import etree
        nsmap = 'http://schemas.openxmlformats.org/drawingml/2006/chart'
        dLbls = chart.series[i]._element.find(f'{{{nsmap}}}dLbls')
        if dLbls is not None:
            sv = dLbls.find(f'{{{nsmap}}}showVal')
            if sv is not None:
                sv.set('val', '1')

    if y_max:
        chart.value_axis.maximum_scale = y_max

    chart.value_axis.tick_labels.font.size = Pt(8)
    chart.value_axis.tick_labels.font.color.rgb = GRAY
    chart.category_axis.tick_labels.font.size = Pt(8)
    chart.category_axis.tick_labels.font.color.rgb = GRAY

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(9)
    chart.has_title = False


def add_pie_chart(slide, x, y, w, h, labels, values, colours, title):
    cd = ChartData()
    cd.categories = labels
    cd.add_series("Summary", values)

    gf = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE,
        Inches(x), Inches(y), Inches(w), Inches(h), cd
    )
    chart = gf.chart

    for i, col in enumerate(colours):
        pt = chart.series[0].points[i]
        pt.format.fill.solid()
        pt.format.fill.fore_color.rgb = col

    dl = chart.series[0].data_labels
    dl.showPercentage = True
    dl.showValue      = False
    dl.font.size      = Pt(11)
    dl.font.color.rgb = RGBColor(0x33, 0x33, 0x33)

    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(9)

    chart.has_title = True
    chart.chart_title.text_frame.text = title
    tf = chart.chart_title.text_frame
    tf.paragraphs[0].font.size = Pt(10)
    tf.paragraphs[0].font.bold = True
    tf.paragraphs[0].font.color.rgb = RGBColor(0x33, 0x33, 0x33)


# ══════════════════════════════════════════════════════════════════════════════
# DATA EXTRACTION
# ══════════════════════════════════════════════════════════════════════════════

def extract_report_data(xlsx_path):
    wb = openpyxl.load_workbook(xlsx_path, data_only=True)
    ws = wb['Sheet1']

    data_sheet = next((s for s in wb.sheetnames if s != 'Sheet1'), None)
    proj_name  = "Project"
    if data_sheet:
        m = re.match(r'ALL_(.+?)_Ratings_List', data_sheet)
        proj_name = m.group(1) if m else data_sheet

    cat_labels, cat_total, cat_issues = [], [], []
    for r in range(3, 10):
        lbl = ws.cell(r, 10).value
        if lbl is None:
            continue
        cat_labels.append(str(lbl))
        cat_total.append(int(ws.cell(r, 11).value or 0))
        cat_issues.append(int(ws.cell(r, 12).value or 0))

    satisfactory = int(ws['Q3'].value or 0)
    observations = int(ws['R3'].value or 0)

    div_labels, div_total, div_issues = [], [], []
    r = 201
    while True:
        asset = ws.cell(r, 1).value
        if asset is None:
            break
        div_labels.append(str(asset))
        div_total.append(int(ws.cell(r, 2).value or 0))
        div_issues.append(int(ws.cell(r, 3).value or 0))
        r += 1

    return {
        "proj_name":    proj_name,
        "cat_labels":   cat_labels,
        "cat_total":    cat_total,
        "cat_issues":   cat_issues,
        "satisfactory": satisfactory,
        "observations": observations,
        "div_labels":   div_labels,
        "div_total":    div_total,
        "div_issues":   div_issues,
    }


# ══════════════════════════════════════════════════════════════════════════════
# SLIDE BUILDERS
# ══════════════════════════════════════════════════════════════════════════════

def build_slide1_category(prs, d, slide_num):
    """Slide 1 — Dashboard layout matching the HTML template:
       Header bar | 3 KPI cards (top) | Category bar chart (bottom-left) | Donut/pie (bottom-right)
    """
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = RGBColor(0xF4, 0xF7, 0xF6)  # light grey bg

    total   = d["satisfactory"] + d["observations"]
    sat_pct = round(d["satisfactory"] / total * 100, 1) if total else 0
    obs_pct = round(d["observations"] / total * 100, 1) if total else 0

    # ── Header bar ────────────────────────────────────────
    add_rect(slide, 0, 0, 10, 0.7, RGBColor(0x2C, 0x3E, 0x50))
    add_textbox(slide, 0.2, 0.08, 1.5, 0.54, "HiRATE",
                22, bold=True, colour=WHITE)
    add_textbox(slide, 0.2, 0.42, 1.5, 0.25, "AUDIT REPORT",
                6, bold=False, colour=RGBColor(0xAA, 0xBB, 0xCC))
    add_textbox(slide, 1.9, 0.1, 6.0, 0.5,
                f"HiRATE Observations — {d['proj_name']}",
                16, bold=True, colour=WHITE)
    # Confidential tag
    add_rect(slide, 8.3, 0.17, 1.55, 0.36, RGBColor(0xFF, 0xF3, 0xCD),
             RGBColor(0xFF, 0xEE, 0xBA), 1)
    add_textbox(slide, 8.3, 0.17, 1.55, 0.36, "CONFIDENTIAL",
                7, bold=True, colour=RGBColor(0x85, 0x64, 0x04),
                align=PP_ALIGN.CENTER)

    # ── 3 KPI Cards ───────────────────────────────────────
    card_y  = 0.82
    card_h  = 1.0
    card_bg = WHITE
    CARD_BORDER = RGBColor(0xE0, 0xE8, 0xF0)

    # Card 1 — Total Observations
    add_rect(slide, 0.2, card_y, 2.85, card_h, card_bg, CARD_BORDER, 0.5)
    add_textbox(slide, 0.35, card_y+0.08, 2.5, 0.22, "TOTAL OBSERVATIONS",
                7, bold=True, colour=RGBColor(0x7F, 0x8C, 0x8D))
    add_textbox(slide, 0.35, card_y+0.3, 2.5, 0.42, str(total),
                26, bold=True, colour=RGBColor(0x2C, 0x3E, 0x50))
    add_textbox(slide, 0.35, card_y+0.74, 2.5, 0.2, "Audits Conducted",
                8, colour=RGBColor(0x7F, 0x8C, 0x8D))

    # Card 2 — Satisfactory (green top border)
    add_rect(slide, 3.25, card_y, 2.85, card_h, card_bg, CARD_BORDER, 0.5)
    add_rect(slide, 3.25, card_y, 2.85, 0.045, RGBColor(0x27, 0xAE, 0x60))  # green top stripe
    add_textbox(slide, 3.4, card_y+0.08, 2.5, 0.22, "SATISFACTORY",
                7, bold=True, colour=RGBColor(0x7F, 0x8C, 0x8D))
    add_textbox(slide, 3.4, card_y+0.3, 2.5, 0.42, f"{sat_pct}%",
                26, bold=True, colour=RGBColor(0x27, 0xAE, 0x60))
    add_textbox(slide, 3.4, card_y+0.74, 2.5, 0.2,
                f"{d['satisfactory']:,} Observations",
                8, colour=RGBColor(0x7F, 0x8C, 0x8D))

    # Card 3 — Issues Found (red top border)
    add_rect(slide, 6.3, card_y, 2.85, card_h, card_bg, CARD_BORDER, 0.5)
    add_rect(slide, 6.3, card_y, 2.85, 0.045, RGBColor(0xC0, 0x39, 0x2B))   # red top stripe
    add_textbox(slide, 6.45, card_y+0.08, 2.5, 0.22, "ISSUES FOUND",
                7, bold=True, colour=RGBColor(0x7F, 0x8C, 0x8D))
    add_textbox(slide, 6.45, card_y+0.3, 2.5, 0.42, f"{obs_pct}%",
                26, bold=True, colour=RGBColor(0xC0, 0x39, 0x2B))
    add_textbox(slide, 6.45, card_y+0.74, 2.5, 0.2,
                f"{d['observations']:,} Observations",
                8, colour=RGBColor(0x7F, 0x8C, 0x8D))

    # ── Bottom section: Category chart (left) + Pie (right) ──
    bottom_y = 1.97
    bottom_h = 3.3

    # Category bar chart
    add_bar_chart(
        slide, 0.2, bottom_y, 6.0, bottom_h,
        d["cat_labels"],
        [
            ("Total Audited", d["cat_total"],  NAVY),
            ("No of Issues",  d["cat_issues"], ORANGE),
        ],
        lbl_size=8,
    )

    # Pie/donut chart — white card background
    add_rect(slide, 6.5, bottom_y, 3.3, bottom_h, WHITE, CARD_BORDER, 0.5)
    add_textbox(slide, 6.65, bottom_y+0.12, 3.0, 0.25, "OBSERVATIONS SUMMARY",
                8, bold=True, colour=RGBColor(0x7F, 0x8C, 0x8D))
    add_pie_chart(
        slide, 6.55, bottom_y+0.35, 3.2, 2.3,
        ["Satisfactory", "Observations"],
        [d["satisfactory"], d["observations"]],
        [RGBColor(0x27, 0xAE, 0x60), RGBColor(0xE0, 0xE0, 0xE0)],
        "",   # title handled by textbox above
    )

    # Footer
    add_footer(slide, f"Slide {slide_num}")


def build_slide2_division(prs, d, slide_num):
    """Slide 2 — Division-wise full-width bar chart."""
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    slide.background.fill.solid()
    slide.background.fill.fore_color.rgb = BG_SLD

    add_header(slide, "HiRATE Observations — Division Wise", d["proj_name"])
    add_textbox(slide, 0.3, 0.75, 9.4, 0.15,
                f"HiRATE Observations Division wise — {d['proj_name']}",
                8, italic=True, colour=RGBColor(0x99, 0x99, 0x99))

    div_max  = max(d["div_issues"]) if d["div_issues"] else 1
    div_ymax = (int(div_max * 1.5 / 5) + 1) * 5

    add_bar_chart(
        slide, 0.3, 0.9, 9.4, 4.3,
        d["div_labels"],
        [
            ("Total Audited", d["div_total"],  NAVY),
            ("No of Issues",  d["div_issues"], ORANGE),
        ],
        y_max=div_ymax,
        lbl_size=7,
    )
    add_footer(slide, f"Slide {slide_num}")


# ══════════════════════════════════════════════════════════════════════════════
# BUILD
# ══════════════════════════════════════════════════════════════════════════════

def build_ppt(report_files, output_path):
    projects = []
    for f in report_files:
        print(f"  Reading: {os.path.basename(f)}")
        try:
            d = extract_report_data(f)
            projects.append(d)
            print(f"    → {d['proj_name']}  cats={len(d['cat_labels'])}  divs={len(d['div_labels'])}")
        except Exception as e:
            print(f"    ⚠ Skipped: {e}")

    if not projects:
        print("No valid report files.")
        return False

    prs = Presentation()
    prs.slide_width  = SW
    prs.slide_height = SH

    for i, d in enumerate(projects):
        build_slide1_category(prs, d, i * 2 + 1)
        build_slide2_division(prs, d, i * 2 + 2)

    prs.save(output_path)
    print(f"\n✓  Saved: {output_path}")
    print(f"   {len(projects)} project(s) × 2 slides = {len(projects)*2} slides total")
    return True


def find_report_files(folder):
    return sorted(glob.glob(os.path.join(folder, "*_REPORT.xlsx")))


def main():
    ap = argparse.ArgumentParser(description="HiRATE PPT generator")
    ap.add_argument("files", nargs="*", help="*_REPORT.xlsx files")
    ap.add_argument("--output", "-o", default=None, help="Output .pptx path")
    args = ap.parse_args()

    files = [os.path.abspath(f) for f in args.files] if args.files \
            else find_report_files(SCRIPT_DIR)

    if not files:
        print(f"ERROR: No *_REPORT.xlsx files found in {SCRIPT_DIR}")
        sys.exit(1)

    out = os.path.abspath(args.output) if args.output \
          else os.path.join(SCRIPT_DIR, DEFAULT_OUTPUT)

    print(f"\nHiRATE PPT Generator  ·  {len(files)} file(s)  →  {out}\n{'─'*50}")
    sys.exit(0 if build_ppt(files, out) else 1)


def generate_ppt_from_reports(report_bytes_list, output_path):
    """Called by sipl_app.py. report_bytes_list = [(filename, bytes), ...]"""
    tmp = []
    try:
        for fname, fbytes in report_bytes_list:
            tf = tempfile.NamedTemporaryFile(suffix=".xlsx", delete=False,
                 prefix=os.path.splitext(fname)[0] + "_")
            tf.write(fbytes); tf.close()
            tmp.append(tf.name)
        return build_ppt(tmp, output_path)
    finally:
        for f in tmp:
            try: os.unlink(f)
            except: pass


if __name__ == "__main__":
    main()
